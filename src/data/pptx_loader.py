import os
from typing import List
from langchain.schema import Document
from pptx import Presentation

class PPTXLoader:
    """Loader for PPTX files from local directory"""
    
    def __init__(self, folder_path: str = None):
        if folder_path is None:
            from config.settings import settings
            self.folder_path = settings.PPTX_FOLDER_PATH
        else:
            self.folder_path = folder_path
        
        if not os.path.exists(self.folder_path):
            # # print(f"PPTX folder does not exist: {self.folder_path}")
            self.folder_path = None
    
    def load_pptx_content(self, file_path: str) -> str:
        """Extract text content from PPTX file"""
        try:
            # # print(f"Loading PPTX content from: {file_path}")
            presentation = Presentation(file_path)
            text_runs = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Add slide title
                if slide.shapes.title:
                    slide_text.append(f"Слайд {slide_num + 1}: {slide.shapes.title.text}")
                
                # Add all text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape != slide.shapes.title:  # Avoid duplicate title
                            slide_text.append(shape.text)
                
                if slide_text:
                    text_runs.append("\n".join(slide_text))
            
            content = "\n\n".join(text_runs)
            # # print(f"Extracted {len(content)} characters from {file_path}")
            return content
            
        except Exception as e:
            # print(f"Error loading PPTX file {file_path}: {e}")
            return ""
    
    def load_documents(self, query: str = None) -> List[Document]:
        """Load all PPTX documents from folder"""
        if not self.folder_path or not os.path.exists(self.folder_path):
            # # print(f"PPTX folder not available: {self.folder_path}")
            return []
        
        documents = []
        pptx_files = []
        
        # Find all PPTX files
        for root, dirs, files in os.walk(self.folder_path):
            for file in files:
                if file.lower().endswith('.pptx'):
                    pptx_files.append(os.path.join(root, file))
        
        for i, file_path in enumerate(pptx_files):
            
            content = self.load_pptx_content(file_path)
            if content:
                # Create LangChain Document with metadata
                metadata = {
                    'source': file_path,
                    'title': os.path.basename(file_path),
                    'type': 'pptx',
                    'original_query': query if query else "general",
                    'search_mode': 'local'
                }
                
                document = Document(
                    page_content=content,
                    metadata=metadata
                )
                documents.append(document)
                # # print(f"Successfully loaded PPTX document {i+1}")
            else:
                # # print(f"Failed to load content for PPTX file {i+1}")
                pass
        
        # # print(f"Loaded {len(documents)} PPTX documents")
        return documents
    
    def load_documents_from_query(self, query: str) -> List[Document]:
        """Load PPTX documents for query (currently loads all, could implement filtering)"""
        # For now, load all documents. Could implement content filtering later
        return self.load_documents(query)